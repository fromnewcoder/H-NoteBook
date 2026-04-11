import asyncio
import json
import os
import httpx
from uuid import UUID

from app.workers.celery_app import celery_app
from app.models.export_job import JobStatus, ExportFormat
from app.config import settings


# Persistent event loop for the worker process
_loop = None


def _get_or_create_event_loop():
    """Get or create a persistent event loop for this process."""
    global _loop
    if _loop is None or _loop.is_closed():
        _loop = asyncio.new_event_loop()
        asyncio.set_event_loop(_loop)
    return _loop


@celery_app.task(bind=True)
def run_export_task(self, job_id: str, notebook_id: str, format: str):
    """Run an export job."""
    import asyncio

    async def _export():
        from sqlalchemy import select
        from app.database import async_session_maker
        from app.models.export_job import ExportJob
        from app.models.source import Source, SourceStatus
        from app.services.export_service import update_export_job_status

        async with async_session_maker() as db:
            # Get job
            result = await db.execute(select(ExportJob).where(ExportJob.id == UUID(job_id)))
            job = result.scalar_one_or_none()

            if not job:
                return

            # Update status to processing
            await update_export_job_status(db, UUID(job_id), JobStatus.PROCESSING)

            try:
                # Get all ready sources
                result = await db.execute(
                    select(Source).where(
                        Source.notebook_id == UUID(notebook_id),
                        Source.status == SourceStatus.READY
                    )
                )
                sources = list(result.scalars().all())

                # Concatenate raw content
                combined_content = "\n\n".join(
                    f"--- {s.name} ---\n{s.raw_content or ''}" for s in sources
                )

                # Generate export using MiniMax API
                export_content = await generate_export_content(format, combined_content)

                # Render file
                file_path = await render_export_file(job_id, format, export_content)

                # Update job as done
                await update_export_job_status(
                    db, UUID(job_id), JobStatus.DONE, file_path=file_path
                )

            except Exception as e:
                await update_export_job_status(
                    db, UUID(job_id), JobStatus.FAILED, error_message=str(e)
                )
                raise

    loop = _get_or_create_event_loop()
    loop.run_until_complete(_export())


async def generate_export_content(format: str, content: str) -> str:
    """Generate structured content using MiniMax API."""
    prompts = {
        "pdf": "Generate a structured summary report with title, abstract, and key findings sections based on the following content. Return only the structured content in plain text format:\n\n",
        "mind_map": 'Return JSON for a mind map with nodes and edges structure: {"nodes": [{"id": "1", "label": "Central Topic"}], "edges": [{"from": "1", "to": "2", "label": "related topic"}]}. Based on:\n\n',
        "docx": "Generate content with clear headings and bullet points based on:\n\n",
        "pptx": 'Return JSON array for up to 10 slides: [{"title": "Slide Title", "bullets": ["point 1", "point 2"]}]. Based on:\n\n',
        "xlsx": 'Return JSON with sheets: [{"sheet": "Sheet Name", "headers": ["Column 1", "Column 2"], "rows": [["value1", "value2"]]}]. Based on:\n\n',
    }

    prompt = prompts.get(format, "") + content[:10000]  # Limit content size

    async with httpx.AsyncClient(timeout=300.0) as client:
        payload = {
            "model": settings.minimax_model,
            "messages": [{"role": "user", "content": prompt}],
            "stream": False
        }

        headers = {
            "Authorization": f"Bearer {settings.minimax_api_key}",
            "Content-Type": "application/json"
        }

        response = await client.post(
            f"{settings.minimax_api_base_url}/v1/messages",
            json=payload,
            headers=headers
        )
        response.raise_for_status()

        result = response.json()
        import logging
        logging.warning(f"MiniMax API response: {result}")

        content = result.get("content", None)

        # Handle list of blocks format (MiniMax API standard format)
        if content and isinstance(content, list):
            for block in content:
                if block.get("type") == "text":
                    text = block.get("text", "")
                    # Strip markdown code fences if present
                    text = text.strip()
                    if text.startswith("```json"):
                        text = text[7:]  # Remove ```json prefix
                    if text.startswith("```"):
                        text = text[3:]  # Remove ``` prefix
                    if text.endswith("```"):
                        text = text[:-3]  # Remove ``` suffix
                    return text.strip()

        # Handle direct string content
        if isinstance(content, str):
            return content

        # Log unexpected format for debugging
        logging.warning(f"Unexpected MiniMax API content format: {result}")

        return ""


async def render_export_file(job_id: str, format: str, content: str) -> str:
    """Render the export to a file and return the path."""
    os.makedirs(settings.export_storage_path, exist_ok=True)
    file_path = os.path.join(settings.export_storage_path, f"{job_id}.{format}")

    if format == "pdf":
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        doc = SimpleDocTemplate(file_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        for line in content.split("\n"):
            if line.startswith("# "):
                story.append(Paragraph(f"<b>{line[2:]}</b>", styles["Heading1"]))
            elif line.startswith("## "):
                story.append(Paragraph(f"<b>{line[3:]}</b>", styles["Heading2"]))
            elif line.strip():
                story.append(Paragraph(line, styles["Normal"]))
            story.append(Spacer(1, 12))

        doc.build(story)

    elif format == "mind_map":
        data = json.loads(content)
        from pyvis.network import Network

        net = Network(height="750px", width="100%", directed=True)

        # Add all nodes first
        nodes = data.get("nodes", [{"id": "1", "label": "Topic"}])
        for node in nodes:
            net.add_node(node.get("id", "1"), label=node.get("label", "Topic"))

        # Then add all edges
        for edge in data.get("edges", []):
            net.add_edge(edge.get("from", "1"), edge.get("to", "1"), label=edge.get("label", ""))

        net.save_graph(file_path)

    elif format == "docx":
        from docx import Document

        doc = Document()
        for line in content.split("\n"):
            if line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.strip():
                doc.add_paragraph(line)

        doc.save(file_path)

    elif format == "pptx":
        from pptx import Presentation

        prs = Presentation()
        slides_data = json.loads(content)

        for slide_data in slides_data[:10]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", "Slide")
            bullets = slide_data.get("bullets", [])
            if bullets:
                slide.placeholders[1].text = "\n".join(f"• {b}" for b in bullets)

        prs.save(file_path)

    elif format == "xlsx":
        from openpyxl import Workbook

        wb = Workbook()
        sheets_data = json.loads(content)

        for i, sheet_data in enumerate(sheets_data):
            if i == 0:
                ws = wb.active
                ws.title = sheet_data.get("sheet", "Sheet1")
            else:
                ws = wb.create_sheet(title=sheet_data.get("sheet", f"Sheet{i+1}"))

            headers = sheet_data.get("headers", [])
            for col, header in enumerate(headers, 1):
                ws.cell(row=1, column=col, value=header)

            rows = sheet_data.get("rows", [])
            for row_idx, row in enumerate(rows, 2):
                for col, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col, value=value)

        wb.save(file_path)

    return file_path
